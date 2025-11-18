"""
Document extraction service for converting PDF, Excel, DOCX, PPTX files to JSON format.
Extracts all data and structure from documents.
"""

import json
import logging
from typing import Dict, Any, Optional
from pathlib import Path
import tempfile
import io

try:
    from doctr.io import DocumentFile
    from doctr.models import ocr_predictor
    HAS_DOCTR = True
except ImportError:
    HAS_DOCTR = False

try:
    import pdf2image
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import pytesseract
    from PIL import Image
    HAS_PYTESSERACT = True
except Exception:
    HAS_PYTESSERACT = False

try:
    from docx import Document as DocxDocument
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False

try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

logger = logging.getLogger(__name__)


class DocumentExtractionService:
    """Service to extract data from various document formats and convert to JSON"""

    def __init__(self):
        self.ocr_model = None

    def get_ocr_model(self):
        """Lazy load OCR model to avoid memory overhead if not needed"""
        if not HAS_DOCTR:
            logger.warning("doctr not installed. OCR functionality will be limited.")
            return None
        if self.ocr_model is None:
            try:
                self.ocr_model = ocr_predictor(pretrained=True)
                logger.info("OCR model loaded successfully")
            except Exception as e:
                logger.error(f"Failed to load OCR model: {e}")
                return None
        return self.ocr_model

    def extract_from_file(self, file_path: str) -> Dict[str, Any]:
        """
        Extract data from a file based on its extension.
        Returns a dictionary with extracted data.
        Uses doctr OCR as primary method for best accuracy.
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        logger.info(f"Extracting data from file: {file_path.name} (type: {extension})")

        if extension == ".pdf":
            return self.extract_from_pdf(str(file_path))
        elif extension in [".xlsx", ".xls"]:
            return self.extract_from_excel(str(file_path))
        elif extension == ".docx":
            return self.extract_from_docx(str(file_path))
        elif extension == ".pptx":
            return self.extract_from_pptx(str(file_path))
        elif extension == ".txt":
            return self.extract_from_text(str(file_path))
        elif extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]:
            # Image files: try doctr first, then pytesseract fallback
            return self.extract_from_image(str(file_path))
        else:
            logger.warning(f"Unsupported file type: {extension}")
            return {
                "status": "unsupported",
                "message": f"File type {extension} is not supported",
                "filename": file_path.name,
                "file_type": extension
            }

    def extract_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and metadata from PDF using doctr (primary OCR).
        Falls back to pytesseract only if doctr is unavailable or fails.
        """
        try:
            result = {
                "file_type": "pdf",
                "filename": Path(file_path).name,
                "pages": [],
                "full_text": "",
                "metadata": {},
                "extraction_method": "none"
            }

            # PRIMARY: Use doctr OCR for best accuracy
            if HAS_DOCTR:
                try:
                    logger.info(f"Using doctr OCR for PDF: {Path(file_path).name}")
                    ocr_model = self.get_ocr_model()
                    if ocr_model:
                        doc = DocumentFile.from_pdf(file_path)
                        ocr_result = ocr_model(doc)
                        
                        full_text = ""
                        for page_idx, page in enumerate(ocr_result.pages):
                            page_data = {
                                "page_number": page_idx + 1,
                                "text": "",
                                "blocks": []
                            }
                            
                            page_text = ""
                            for block in page.blocks:
                                block_text = ""
                                block_confidence = 0.0
                                word_count = 0
                                
                                for line in block.lines:
                                    for word in line.words:
                                        block_text += word.value + " "
                                        block_confidence += word.confidence if hasattr(word, 'confidence') else 1.0
                                        word_count += 1
                                
                                if block_text.strip():
                                    avg_confidence = block_confidence / word_count if word_count > 0 else 0.0
                                    page_data["blocks"].append({
                                        "text": block_text.strip(),
                                        "confidence": float(avg_confidence)
                                    })
                                    page_text += block_text + " "
                            
                            page_data["text"] = page_text.strip()
                            full_text += page_text + "\n"
                            result["pages"].append(page_data)
                        
                        result["full_text"] = full_text.strip()
                        result["status"] = "success"
                        result["extraction_method"] = "doctr_ocr"
                        logger.info(f"✓ Successfully extracted PDF using doctr: {Path(file_path).name}")
                        return result
                except Exception as e:
                    logger.error(f"doctr OCR extraction failed: {e}")
                    # Continue to fallback

            # FALLBACK 1: pytesseract (if doctr not available or failed)
            if HAS_PYTESSERACT and HAS_PDF2IMAGE:
                try:
                    logger.info(f"Falling back to pytesseract for PDF: {Path(file_path).name}")
                    pages_text = []
                    images = pdf2image.convert_from_path(file_path)
                    
                    for i, img in enumerate(images):
                        text = pytesseract.image_to_string(img)
                        pages_text.append(text)
                    
                    full_text = "\n".join(pages_text).strip()
                    result["pages"] = [
                        {"page_number": i + 1, "text": p} for i, p in enumerate(pages_text)
                    ]
                    result["full_text"] = full_text
                    result["status"] = "success"
                    result["extraction_method"] = "pytesseract_fallback"
                    logger.info(f"✓ Successfully extracted PDF using pytesseract: {Path(file_path).name}")
                    return result
                except Exception as e:
                    logger.error(f"pytesseract PDF fallback failed: {e}")

            # FALLBACK 2: pdf2image (basic image info, no text)
            if HAS_PDF2IMAGE:
                try:
                    logger.info(f"Falling back to basic pdf2image conversion: {Path(file_path).name}")
                    images = pdf2image.convert_from_path(file_path)
                    result["pages"] = [
                        {
                            "page_number": i + 1,
                            "has_image": True,
                            "image_size": f"{img.width}x{img.height}"
                        }
                        for i, img in enumerate(images)
                    ]
                    result["status"] = "partial"
                    result["extraction_method"] = "pdf2image_only"
                    result["message"] = "PDF converted to images but no OCR engine available. Install doctr or enable pytesseract."
                    logger.warning(f"PDF extracted as images only (no OCR): {Path(file_path).name}")
                    return result
                except Exception as e:
                    logger.error(f"pdf2image fallback failed: {e}")

            # No OCR method available
            result["status"] = "error"
            result["message"] = "No OCR method available. Install python-doctr or pytesseract."
            logger.error(f"Cannot extract PDF: {Path(file_path).name} - No OCR method available")
            return result

        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            return {
                "file_type": "pdf",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def extract_from_excel(self, file_path: str) -> Dict[str, Any]:
        """Extract data from Excel files"""
        try:
            result = {
                "file_type": "excel",
                "filename": Path(file_path).name,
                "sheets": [],
                "status": "success"
            }

            # Try pandas first (better for structured data)
            if HAS_PANDAS:
                try:
                    excel_file = pd.ExcelFile(file_path)
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        result["sheets"].append({
                            "name": sheet_name,
                            "rows": len(df),
                            "columns": list(df.columns),
                            "data": df.fillna("").to_dict(orient="records")
                        })
                    result["extraction_method"] = "pandas"
                    return result
                except Exception as e:
                    logger.warning(f"Pandas extraction failed, trying openpyxl: {e}")

            # Fallback to openpyxl
            if HAS_OPENPYXL:
                try:
                    workbook = load_workbook(file_path)
                    for sheet_name in workbook.sheetnames:
                        worksheet = workbook[sheet_name]
                        rows = []
                        headers = None
                        
                        for row_idx, row in enumerate(worksheet.iter_rows(values_only=True), 1):
                            if row_idx == 1:
                                headers = [str(cell) if cell else f"Column_{i}" for i, cell in enumerate(row)]
                            else:
                                if headers:
                                    row_dict = {headers[i]: cell for i, cell in enumerate(row)}
                                    rows.append(row_dict)
                        
                        result["sheets"].append({
                            "name": sheet_name,
                            "rows": len(rows),
                            "columns": headers or [],
                            "data": rows
                        })
                    
                    result["extraction_method"] = "openpyxl"
                    return result
                except Exception as e:
                    logger.error(f"openpyxl extraction failed: {e}")

            result["status"] = "error"
            result["message"] = "Could not extract Excel. Install 'pandas' or 'openpyxl'."
            return result

        except Exception as e:
            logger.error(f"Error extracting Excel: {e}")
            return {
                "file_type": "excel",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def extract_from_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from DOCX files"""
        try:
            if not HAS_PYTHON_DOCX:
                return {
                    "file_type": "docx",
                    "filename": Path(file_path).name,
                    "status": "error",
                    "message": "python-docx not installed"
                }

            doc = DocxDocument(file_path)
            result = {
                "file_type": "docx",
                "filename": Path(file_path).name,
                "paragraphs": [],
                "tables": [],
                "full_text": "",
                "status": "success"
            }

            # Extract paragraphs
            full_text = ""
            for para in doc.paragraphs:
                if para.text.strip():
                    result["paragraphs"].append({
                        "text": para.text,
                        "level": para.style.name if para.style else "Normal"
                    })
                    full_text += para.text + "\n"

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = {
                    "table_number": table_idx + 1,
                    "rows": []
                }
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data["rows"].append(row_data)
                result["tables"].append(table_data)

            result["full_text"] = full_text.strip()
            return result

        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            return {
                "file_type": "docx",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def extract_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract text and structure from PPTX files"""
        try:
            if not HAS_PYTHON_PPTX:
                return {
                    "file_type": "pptx",
                    "filename": Path(file_path).name,
                    "status": "error",
                    "message": "python-pptx not installed"
                }

            prs = Presentation(file_path)
            result = {
                "file_type": "pptx",
                "filename": Path(file_path).name,
                "slides": [],
                "full_text": "",
                "status": "success"
            }

            full_text = ""
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "shapes": []
                }

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["shapes"].append({
                            "type": shape.shape_type,
                            "text": shape.text
                        })
                        full_text += shape.text + "\n"

                result["slides"].append(slide_data)

            result["full_text"] = full_text.strip()
            return result

        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return {
                "file_type": "pptx",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def extract_from_text(self, file_path: str) -> Dict[str, Any]:
        """Extract data from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            return {
                "file_type": "text",
                "filename": Path(file_path).name,
                "content": content,
                "lines": content.split('\n'),
                "status": "success"
            }

        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            return {
                "file_type": "text",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def extract_from_image(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from image files using doctr (primary) or pytesseract (fallback).
        Uses doctr for best accuracy when available.
        """
        try:
            result = {
                "file_type": "image",
                "filename": Path(file_path).name,
                "full_text": "",
                "status": "error"
            }

            # PRIMARY: Use doctr OCR for best accuracy
            if HAS_DOCTR:
                try:
                    logger.info(f"Using doctr OCR for image: {Path(file_path).name}")
                    ocr_model = self.get_ocr_model()
                    if ocr_model:
                        doc = DocumentFile.from_images(file_path)
                        ocr_result = ocr_model(doc)
                        
                        full_text = ""
                        for page in ocr_result.pages:
                            for block in page.blocks:
                                for line in block.lines:
                                    for word in line.words:
                                        full_text += word.value + " "
                        
                        result["full_text"] = full_text.strip()
                        result["status"] = "success"
                        result["extraction_method"] = "doctr_ocr"
                        logger.info(f"✓ Successfully extracted image using doctr: {Path(file_path).name}")
                        return result
                except Exception as e:
                    logger.error(f"doctr OCR extraction failed for image: {e}")
                    # Continue to fallback

            # FALLBACK: pytesseract
            if HAS_PYTESSERACT:
                try:
                    logger.info(f"Falling back to pytesseract for image: {Path(file_path).name}")
                    img = Image.open(file_path)
                    text = pytesseract.image_to_string(img)
                    result["full_text"] = text.strip()
                    result["status"] = "success"
                    result["extraction_method"] = "pytesseract_fallback"
                    logger.info(f"✓ Successfully extracted image using pytesseract: {Path(file_path).name}")
                    return result
                except Exception as e:
                    logger.error(f"pytesseract image fallback failed: {e}")

            # No OCR method available
            result["status"] = "error"
            result["message"] = "No OCR method available. Install python-doctr or pytesseract."
            logger.error(f"Cannot extract image: {Path(file_path).name} - No OCR method available")
            return result

        except Exception as e:
            logger.error(f"Error extracting image: {e}")
            return {
                "file_type": "image",
                "filename": Path(file_path).name,
                "status": "error",
                "error": str(e)
            }

    def ocr_image_fallback(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Deprecated: Use extract_from_image instead. Fallback OCR for image files using pytesseract."""
        # This method is kept for backward compatibility but is no longer the primary path
        return self.extract_from_image(file_path)


# Create a singleton instance
extraction_service = DocumentExtractionService()
