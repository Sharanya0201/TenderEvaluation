import io
import asyncio
import os
import tempfile
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from typing import Dict, List, Optional
import logging

from doctr.io import DocumentFile
from doctr.models import ocr_predictor
from sqlalchemy.orm import Session

# Import conversion libraries
import pdf2image
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

from app.models.user import VendorDocument, OCRResult

logger = logging.getLogger(__name__)

class OCRService:
    def __init__(self):
        self.ocr_model = None
        self.executor = ThreadPoolExecutor(max_workers=2)
    
    def get_ocr_model(self):
        """Lazy load OCR model"""
        if self.ocr_model is None:
            try:
                self.ocr_model = ocr_predictor(pretrained=True)
                logger.info("OCR model loaded successfully")
            except Exception as e:
                logger.error(f"Failed to load OCR model: {e}")
                raise
        return self.ocr_model
    
    def get_or_create_ocr_result(self, db: Session, document_id: int) -> OCRResult:
        """Get existing OCR result or create a new one"""
        ocr_result = db.query(OCRResult).filter(
            OCRResult.document_id == document_id
        ).first()
        
        if not ocr_result:
            ocr_result = OCRResult(
                document_id=document_id,
                status='pending'
            )
            db.add(ocr_result)
            db.commit()
            db.refresh(ocr_result)
        
        return ocr_result
    
    def convert_docx_to_text(self, file_content: bytes) -> str:
        """Convert DOCX to text"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            doc = DocxDocument(temp_file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Clean up
            os.unlink(temp_file_path)
            
            return text_content
            
        except Exception as e:
            logger.error(f"DOCX conversion failed: {e}")
            raise
    
    def convert_pptx_to_text(self, file_content: bytes) -> str:
        """Convert PPTX to text"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            prs = Presentation(temp_file_path)
            text_content = ""
            
            for i, slide in enumerate(prs.slides):
                text_content += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content += shape.text + "\n"
                text_content += "\n"
            
            # Clean up
            os.unlink(temp_file_path)
            
            return text_content
            
        except Exception as e:
            logger.error(f"PPTX conversion failed: {e}")
            raise
    
    def convert_xlsx_to_text(self, file_content: bytes) -> str:
        """Convert XLSX to text"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            wb = load_workbook(temp_file_path)
            text_content = ""
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content += f"=== Sheet: {sheet_name} ===\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    text_content += row_text + "\n"
                text_content += "\n"
            
            # Clean up
            os.unlink(temp_file_path)
            
            return text_content
            
        except Exception as e:
            logger.error(f"XLSX conversion failed: {e}")
            raise
    
    def extract_text_from_file(self, file_content: bytes, filename: str, content_type: str) -> tuple[str, float]:
        """Extract text from file using appropriate method"""
        try:
            file_extension = filename.lower().split('.')[-1] if '.' in filename else ''
            
            # Handle different file formats
            if file_extension in ['doc', 'docx']:
                logger.info(f"Converting DOCX to text: {filename}")
                text_content = self.convert_docx_to_text(file_content)
                return text_content, 1.0
                
            elif file_extension in ['ppt', 'pptx']:
                logger.info(f"Converting PPTX to text: {filename}")
                text_content = self.convert_pptx_to_text(file_content)
                return text_content, 1.0
                
            elif file_extension in ['xls', 'xlsx']:
                logger.info(f"Converting XLSX to text: {filename}")
                text_content = self.convert_xlsx_to_text(file_content)
                return text_content, 1.0
                
            elif file_extension in ['txt', 'csv']:
                logger.info(f"Reading text file directly: {filename}")
                text_content = file_content.decode('utf-8', errors='ignore')
                return text_content, 1.0
                
            else:
                # Use OCR for images, PDFs, and other formats
                logger.info(f"Using OCR for file: {filename}")
                return self.extract_text_with_ocr(file_content, filename)
                
        except Exception as e:
            logger.error(f"Text extraction failed for {filename}: {e}")
            raise
        
    def extract_text_with_ocr(self, file_content: bytes, filename: str) -> tuple[str, float]:
        """Extract text from file using OCR"""
        try:
            print(f"🔍 Starting OCR for file: {filename}")
            print(f"📁 File size: {len(file_content)} bytes")
            print(f"📊 First 100 bytes: {file_content[:100]}")
            print(f"🔢 Is file content empty? {len(file_content) == 0}")
            
            # Check if file content is actually a PDF
            if filename.lower().endswith('.pdf'):
                print(f"📄 File starts with PDF header? {file_content.startswith(b'%PDF')}")
                if not file_content.startswith(b'%PDF'):
                    print("❌ WARNING: File doesn't have PDF header!")
            
            # Check if we can actually read the file
            try:
                file_bytes = io.BytesIO(file_content)
                print(f"✅ Successfully created BytesIO object")
            except Exception as e:
                print(f"❌ Failed to create BytesIO: {e}")
                return "", 0.0
            
            model = self.get_ocr_model()
            print(f"🤖 OCR model loaded: {model is not None}")
            
            try:
                if filename.lower().endswith('.pdf'):
                    print("📖 Loading as PDF...")
                    doc = DocumentFile.from_pdf(file_bytes)
                    print(f"✅ PDF loaded, pages: {len(doc)}")
                else:
                    print("🖼️ Loading as image...")
                    doc = DocumentFile.from_images(file_bytes)
                    print(f"✅ Image loaded")
            except Exception as e:
                print(f"❌ Failed to load document: {e}")
                return "", 0.0
            
            # Run OCR
            print("🚀 Running OCR...")
            result = model(doc)
            print(f"✅ OCR completed")
            print(f"📄 Pages detected: {len(result.pages)}")
            
            # Extract text and confidence
            ocr_text = ""
            total_confidence = 0.0
            word_count = 0
            
            for page_num, page in enumerate(result.pages):
                print(f"📖 Processing page {page_num + 1}")
                print(f"   Blocks on page: {len(page.blocks)}")
                
                for block_num, block in enumerate(page.blocks):
                    print(f"   📦 Block {block_num + 1}, lines: {len(block.lines)}")
                    
                    for line_num, line in enumerate(block.lines):
                        print(f"      📝 Line {line_num + 1}, words: {len(line.words)}")
                        
                        for word_num, word in enumerate(line.words):
                            ocr_text += word.value + " "
                            total_confidence += word.confidence
                            word_count += 1
                            if word_num < 3:  # Show first 3 words for debugging
                                print(f"         🔤 Word {word_num + 1}: '{word.value}' (confidence: {word.confidence:.2f})")
            
            avg_confidence = total_confidence / word_count if word_count > 0 else 0.0
            print(f"📊 FINAL RESULTS:")
            print(f"   Words found: {word_count}")
            print(f"   Average confidence: {avg_confidence:.2f}")
            print(f"   Text sample: {ocr_text[:100]}...")
            print(f"   Total text length: {len(ocr_text)}")
            
            return ocr_text.strip(), avg_confidence
            
        except Exception as e:
            print(f"💥 OCR extraction failed for {filename}: {e}")
            import traceback
            traceback.print_exc()
            raise
        
    async def process_document_ocr(self, db: Session, document_id: int) -> Dict:
        """Process OCR for a single document"""
        # Check if document exists
        document = db.query(VendorDocument).filter(VendorDocument.id == document_id).first()
        if not document:
            raise ValueError("Document not found")
        
        # Get or create OCR result
        ocr_result = self.get_or_create_ocr_result(db, document_id)
        
        # Check if already processing
        if ocr_result.status == 'processing':
            raise ValueError("OCR processing already in progress")
        
        # Update OCR result status
        ocr_result.status = 'processing'
        ocr_result.processed_at = None
        ocr_result.error_message = None
        db.commit()
        
        try:
            # Process in background thread
            def process_ocr_task():
                try:
                    # Extract text using appropriate method
                    ocr_text, confidence = self.extract_text_from_file(
                        document.file_content,
                        document.original_filename,
                        document.content_type or ''
                    )
                    
                    # Update OCR result
                    ocr_result.ocr_text = ocr_text
                    ocr_result.confidence = confidence
                    ocr_result.status = 'completed'
                    ocr_result.processed_at = datetime.utcnow()
                    ocr_result.error_message = None
                    db.commit()
                    
                    logger.info(f"OCR completed for document {document_id}: {document.original_filename}")
                    
                except Exception as e:
                    logger.error(f"OCR processing failed for document {document_id}: {e}")
                    ocr_result.status = 'failed'
                    ocr_result.error_message = str(e)
                    ocr_result.processed_at = datetime.utcnow()
                    db.commit()
                    raise
            
            # Run in thread pool
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(self.executor, process_ocr_task)
            
            return {
                'document_id': document_id,
                'status': 'processing'
            }
            
        except Exception as e:
            ocr_result.status = 'failed'
            ocr_result.error_message = str(e)
            ocr_result.processed_at = datetime.utcnow()
            db.commit()
            raise
    
    async def bulk_process_ocr(self, db: Session, document_ids: List[int]) -> Dict:
        """Bulk OCR processing for multiple documents"""
        batch_id = f"batch_{int(datetime.utcnow().timestamp())}"
        
        async def process_batch():
            for doc_id in document_ids:
                try:
                    await self.process_document_ocr(db, doc_id)
                except Exception as e:
                    logger.error(f"Failed to process document {doc_id}: {e}")
        
        # Start batch processing in background
        asyncio.create_task(process_batch())
        
        return {
            'batch_id': batch_id,
            'total_documents': len(document_ids),
            'processed_documents': 0,
            'status': 'processing'
        }
    
    def get_ocr_status(self, db: Session, document_id: int) -> Dict:
        """Get OCR status for a document"""
        ocr_result = db.query(OCRResult).filter(
            OCRResult.document_id == document_id
        ).first()
        
        if not ocr_result:
            return {
                'status': 'pending',
                'ocr_text': None,
                'confidence': None,
                'processed_at': None,
                'error_message': None
            }
        
        return {
            'status': ocr_result.status,
            'ocr_text': ocr_result.ocr_text or ocr_result.corrected_text,
            'confidence': ocr_result.confidence,
            'processed_at': ocr_result.processed_at,
            'error_message': ocr_result.error_message
        }
    
    def correct_ocr_text(self, db: Session, document_id: int, corrected_text: str) -> bool:
        """Manual OCR text correction"""
        ocr_result = db.query(OCRResult).filter(
            OCRResult.document_id == document_id
        ).first()
        
        if not ocr_result:
            # Create new OCR result if doesn't exist
            ocr_result = self.get_or_create_ocr_result(db, document_id)
        
        # Update corrected text
        ocr_result.corrected_text = corrected_text
        ocr_result.status = 'corrected'
        db.commit()
        
        return True

# Global OCR service instance
ocr_service = OCRService()